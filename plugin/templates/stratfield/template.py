"""
Stratfield PowerPoint Template Generator

This module provides a presentation template using the Stratfield corporate template.

Usage:
    from templates import get_template

    prs = get_template("stratfield")
    prs.add_title_slide("My Presentation", "A great subtitle", "January 2025")
    prs.add_section_break("Section 1")
    prs.add_content_slide("Slide Title", "", [
        ("First bullet point", 0),
        ("Sub-bullet", 1),
        ("Another point", 0),
    ])
    prs.save("output.pptx")
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from plugin.lib.presentation.template_base import PresentationTemplate
from plugin.templates import register_template


# ============================================================================
# Constants
# ============================================================================

# Slide dimensions (16:9)
SLIDE_WIDTH_EMU = 9144000
SLIDE_HEIGHT_EMU = 5143500

# EMU conversion
EMU_PER_INCH = 914400
EMU_PER_PT = 12700


# ============================================================================
# Color Utilities
# ============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


# Brand color palette
COLORS = {
    # Primary
    'primary_green': '#00764F',
    'dark_teal': '#296057',
    'section_bg': '#30555E',
    'bg_light': '#F2F2F2',

    # Text
    'text_dark': '#637878',
    'text_dark_alt': '#32645F',
    'text_light': '#6E787B',
    'title_cream': '#DDFAD3',
    'title_cream_alt': '#DDF9D3',
    'subtitle_light': '#F2F2F2',
    'white': '#FFFFFF',

    # Accents
    'accent_teal': '#82B2B4',
    'accent_lime': '#BCC365',
    'accent_gold': '#DDB945',
    'accent_orange': '#DD9624',
    'accent_red': '#A91E2A',
}


# ============================================================================
# Font Configuration
# ============================================================================

FONTS = {
    'title_black': 'Avenir Black',
    'title_heavy': 'Avenir Heavy',
    'body': 'Avenir',
    'body_medium': 'Avenir Medium',
    'body_light': 'Avenir Light',
}

# Fallbacks if Avenir not available
FONT_FALLBACKS = {
    'Avenir Black': 'Arial Black',
    'Avenir Heavy': 'Arial Bold',
    'Avenir': 'Arial',
    'Avenir Medium': 'Arial',
    'Avenir Light': 'Arial',
}


# ============================================================================
# Bullet Configuration (matches original template)
# ============================================================================

BULLET_SPECS = {
    # Level 0: Main bullets
    0: {
        'marL': 170089,      # Left margin (~0.186")
        'indent': -170089,   # Hanging indent
        'size': 14,
        'space_before': 8,
    },
    # Level 1: Sub-bullets
    1: {
        'marL': 340688,      # Left margin (~0.373")
        'indent': -109728,   # Hanging indent (~0.120")
        'size': 12,
        'space_before': 4,
    },
    # Level 2: Sub-sub-bullets
    2: {
        'marL': 572516,      # Left margin (~0.627")
        'indent': -170089,   # Hanging indent
        'size': 11,
        'space_before': 4,
    },
}

# Bullet styling
BULLET_CHAR = '§'           # Small square in Wingdings
BULLET_FONT = 'Wingdings'
BULLET_COLOR = '00764F'     # Primary green (accent1)


# ============================================================================
# Layout Specifications (in inches for readability)
# ============================================================================

LAYOUTS = {
    'title_slide': {
        'background_color': None,  # Uses background image
        'logo': {'x': 0.40, 'y': 0.97, 'w': 2.33, 'h': 1.21},
        'title': {
            'x': 0.54, 'y': 2.40, 'w': 8.44, 'h': 0.60,
            'font': 'Avenir Black', 'size': 32, 'color': '#DDFAD3',
            'bold': True, 'align': 'left'
        },
        'subtitle': {
            'x': 0.61, 'y': 3.11, 'w': 8.56, 'h': 0.44,
            'font': 'Avenir', 'size': 20, 'color': '#F2F2F2',
            'align': 'left'
        },
        'date': {
            'x': 0.61, 'y': 3.96, 'w': 2.47, 'h': 0.37,
            'font': 'Avenir Light', 'size': 16, 'color': '#F2F2F2',
            'align': 'left'
        },
    },

    'section_break': {
        'background_color': '#30555E',
        'title': {
            'x': 0.0, 'y': 1.55, 'w': 10.0, 'h': 1.23,
            'font': 'Avenir Black', 'size': 32, 'color': '#DDF9D3',
            'bold': True, 'align': 'center', 'anchor': 'bottom'
        },
        'logo_mark': {'x': 4.38, 'y': 2.81, 'w': 1.25, 'h': 1.96},
        'slide_number': {
            'x': 9.57, 'y': 5.44, 'w': 0.41, 'h': 0.20,
            'font': 'Avenir Medium', 'size': 8, 'color': '#FFFFFF',
            'align': 'right', 'anchor': 'center'
        },
    },

    'content_text_only': {
        'background_color': '#F2F2F2',
        'title': {
            'x': 0.16, 'y': 0.07, 'w': 8.82, 'h': 0.41,
            'font': 'Avenir Black', 'size': 20, 'color': '#296057',
            'align': 'left', 'anchor': 'center'
        },
        'body': {
            'x': 0.28, 'y': 0.72, 'w': 9.39, 'h': 4.51,
            'font': 'Avenir', 'color': '#637878',
            'levels': {
                0: {'size': 14, 'space_before': 8},
                1: {'size': 12, 'space_before': 4},
                2: {'size': 11, 'space_before': 4},
            }
        },
        'footer_bar': {'x': 0.0, 'y': 5.33, 'w': 10.0, 'h': 0.31},
        'footer_logo': {'x': 0.0, 'y': 5.32, 'w': 0.79, 'h': 0.31},
        'slide_number': {
            'x': 9.57, 'y': 5.44, 'w': 0.41, 'h': 0.20,
            'font': 'Avenir Medium', 'size': 8, 'color': '#FFFFFF',
            'align': 'right', 'anchor': 'center'
        },
    },

    'content_full_image': {
        'background_color': '#F2F2F2',
        'title': {
            'x': 0.16, 'y': 0.07, 'w': 8.82, 'h': 0.41,
            'font': 'Avenir Black', 'size': 20, 'color': '#296057',
            'align': 'left', 'anchor': 'center'
        },
        'image': {
            'x': 0.40, 'y': 0.66, 'w': 9.20, 'h': 4.30,
        },
        'footer_bar': {'x': 0.0, 'y': 5.33, 'w': 10.0, 'h': 0.31},
        'footer_logo': {'x': 0.0, 'y': 5.32, 'w': 0.79, 'h': 0.31},
        'slide_number': {
            'x': 9.57, 'y': 5.44, 'w': 0.41, 'h': 0.20,
            'font': 'Avenir Medium', 'size': 8, 'color': '#FFFFFF',
            'align': 'right', 'anchor': 'center'
        },
    },

    'content_text_and_image': {
        'background_color': '#F2F2F2',
        'left_panel': {
            'x': 0.0, 'y': 0.0, 'w': 3.41, 'h': 5.34,
            'fill': '#296057'
        },
        'title': {
            'x': 0.16, 'y': 0.07, 'w': 3.16, 'h': 0.41,
            'font': 'Avenir Black', 'size': 14, 'color': '#FFFFFF',
            'align': 'left'
        },
        'body': {
            'x': 0.18, 'y': 0.63, 'w': 3.16, 'h': 4.62,
            'font': 'Avenir', 'color': '#FFFFFF',
            'levels': {
                0: {'size': 14, 'space_before': 8},
                1: {'size': 12, 'space_before': 4},
                2: {'size': 11, 'space_before': 4},
            }
        },
        'image': {
            'x': 3.85, 'y': 0.44, 'w': 5.70, 'h': 4.50,
        },
        'footer_bar': {'x': 0.0, 'y': 5.33, 'w': 10.0, 'h': 0.31},
        'footer_logo': {'x': 0.0, 'y': 5.32, 'w': 0.79, 'h': 0.31},
        'slide_number': {
            'x': 9.57, 'y': 5.44, 'w': 0.41, 'h': 0.20,
            'font': 'Avenir Medium', 'size': 8, 'color': '#FFFFFF',
            'align': 'right', 'anchor': 'center'
        },
    },
}


# ============================================================================
# Main Presentation Template Class
# ============================================================================

@register_template
class StratfieldPresentation(PresentationTemplate):
    """
    Create PowerPoint presentations using the Stratfield template.

    This template features green and teal professional colors with Avenir font.

    Example:
        from templates import get_template

        prs = get_template("stratfield")
        prs.add_title_slide("Title", "Subtitle", "Date")
        prs.add_content_slide("Slide Title", "", [("Bullet 1", 0), ("Sub", 1)])
        prs.save("output.pptx")
    """

    @property
    def name(self) -> str:
        """Human-readable template name."""
        return "Stratfield Consulting"

    @property
    def id(self) -> str:
        """Template identifier."""
        return "stratfield"

    @property
    def description(self) -> str:
        """Template description."""
        return "Green and teal professional template with Avenir font"

    def __init__(self, assets_dir: str = None):
        """
        Initialize a new Stratfield presentation.

        Args:
            assets_dir: Path to directory containing template assets.
                       Defaults to ./assets relative to this module.
        """
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_WIDTH_EMU)
        self.prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

        if assets_dir:
            self.assets_dir = Path(assets_dir)
        else:
            self.assets_dir = Path(__file__).parent / 'assets'

        self._slide_count = 0

    def _get_blank_layout(self):
        """Get or create a blank slide layout."""
        return self.prs.slide_layouts[6]  # Usually blank

    def _add_textbox(self, slide, spec: dict, text: str):
        """Add a textbox to the slide using specification dict."""
        shape = slide.shapes.add_textbox(
            Inches(spec['x']),
            Inches(spec['y']),
            Inches(spec['w']),
            Inches(spec['h'])
        )

        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text

        if 'font' in spec:
            p.font.name = spec['font']
        if 'size' in spec:
            p.font.size = Pt(spec['size'])
        if 'color' in spec:
            p.font.color.rgb = hex_to_rgb(spec['color'])
        if spec.get('bold'):
            p.font.bold = True

        # Alignment
        align = spec.get('align', 'left')
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        # Vertical anchor
        anchor = spec.get('anchor', 'top')
        if anchor == 'center':
            tf.paragraphs[0].alignment
            shape.text_frame.auto_size = None
        elif anchor == 'bottom':
            pass  # Would need MSO_ANCHOR.BOTTOM

        return shape

    def _add_image(self, slide, spec: dict, image_path: str):
        """Add an image to the slide."""
        if not os.path.exists(image_path):
            # Create placeholder rectangle if image not found
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(spec['x']),
                Inches(spec['y']),
                Inches(spec['w']),
                Inches(spec['h'])
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb('#CCCCCC')
            return shape

        return slide.shapes.add_picture(
            image_path,
            Inches(spec['x']),
            Inches(spec['y']),
            Inches(spec['w']),
            Inches(spec['h'])
        )

    def _add_rectangle(self, slide, spec: dict):
        """Add a filled rectangle shape."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(spec['x']),
            Inches(spec['y']),
            Inches(spec['w']),
            Inches(spec['h'])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(spec['fill'])
        shape.line.fill.background()  # No border
        return shape

    def _apply_bullet_formatting(self, paragraph, level: int, text_color: str = '#637878'):
        """
        Apply PowerPoint native bullet formatting to a paragraph.

        Args:
            paragraph: The paragraph object to format
            level: Bullet level (0, 1, or 2)
            text_color: Hex color for the text
        """
        spec = BULLET_SPECS.get(level, BULLET_SPECS[0])

        # Set text properties first (this creates defRPr element)
        paragraph.font.name = FONTS['body']
        paragraph.font.size = Pt(spec['size'])
        paragraph.font.color.rgb = hex_to_rgb(text_color)
        paragraph.space_before = Pt(spec['space_before'])
        paragraph.level = level

        # Access the underlying XML element for bullet properties
        pPr = paragraph._p.get_or_add_pPr()

        # Set paragraph margins for proper bullet indentation
        pPr.set('marL', str(spec['marL']))
        pPr.set('indent', str(spec['indent']))

        # Remove any existing bullet settings
        for child in list(pPr):
            tag_name = etree.QName(child.tag).localname
            if tag_name in ('buNone', 'buChar', 'buFont', 'buClr', 'buAutoNum', 'buBlip'):
                pPr.remove(child)

        # Find defRPr to insert bullet elements BEFORE it
        # PowerPoint requires strict element ordering: buClr, buFont, buChar must come before defRPr
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is not None:
            insert_index = list(pPr).index(defRPr)
        else:
            insert_index = len(list(pPr))

        # Create bullet elements
        buClr = etree.Element(qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', BULLET_COLOR)

        buFont = etree.Element(qn('a:buFont'))
        buFont.set('typeface', BULLET_FONT)

        buChar = etree.Element(qn('a:buChar'))
        buChar.set('char', BULLET_CHAR)

        # Insert in correct order: buClr, buFont, buChar (before defRPr)
        pPr.insert(insert_index, buChar)
        pPr.insert(insert_index, buFont)
        pPr.insert(insert_index, buClr)

    def _set_background_color(self, slide, color: str):
        """Set solid background color for slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color)

    def _add_footer_elements(self, slide, layout: dict):
        """Add footer bar (gradient image), logo, and slide number."""
        # Footer bar - use gradient image if available
        if 'footer_bar' in layout:
            fb = layout['footer_bar']
            footer_bar_path = self.assets_dir / 'footer_bar.png'
            if footer_bar_path.exists():
                slide.shapes.add_picture(
                    str(footer_bar_path),
                    Inches(fb['x']), Inches(fb['y']),
                    Inches(fb['w']), Inches(fb['h'])
                )
            else:
                # Fallback to solid color if image not found
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(fb['x']), Inches(fb['y']),
                    Inches(fb['w']), Inches(fb['h'])
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = hex_to_rgb(COLORS['primary_green'])
                bar.line.fill.background()

        # Footer logo
        if 'footer_logo' in layout:
            fl = layout['footer_logo']
            logo_path = self.assets_dir / 'logo_footer.png'
            if logo_path.exists():
                slide.shapes.add_picture(
                    str(logo_path),
                    Inches(fl['x']), Inches(fl['y']),
                    Inches(fl['w']), Inches(fl['h'])
                )

        # Slide number
        if 'slide_number' in layout:
            sn = layout['slide_number']
            shape = slide.shapes.add_textbox(
                Inches(sn['x']), Inches(sn['y']),
                Inches(sn['w']), Inches(sn['h'])
            )
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = str(self._slide_count)
            p.font.name = sn.get('font', 'Avenir Medium')
            p.font.size = Pt(sn.get('size', 8))
            p.font.color.rgb = hex_to_rgb(sn.get('color', '#FFFFFF'))
            p.alignment = PP_ALIGN.RIGHT

    def add_title_slide(self, title: str, subtitle: str = "", date: str = "") -> None:
        """
        Add a title/cover slide.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle
            date: Optional date string
        """
        self._slide_count += 1
        layout = LAYOUTS['title_slide']

        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Background - use gradient image if available, otherwise solid color
        bg_path = self.assets_dir / 'title_background.png'
        if bg_path.exists():
            slide.shapes.add_picture(
                str(bg_path),
                Inches(0), Inches(0),
                Inches(10), Inches(5.625)
            )
        else:
            self._set_background_color(slide, COLORS['primary_green'])

        # Logo
        logo_path = self.assets_dir / 'logo_title.png'
        if logo_path.exists() and 'logo' in layout:
            logo_spec = layout['logo']
            slide.shapes.add_picture(
                str(logo_path),
                Inches(logo_spec['x']), Inches(logo_spec['y']),
                Inches(logo_spec['w']), Inches(logo_spec['h'])
            )

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Subtitle
        if subtitle:
            self._add_textbox(slide, layout['subtitle'], subtitle)

        # Date
        if date:
            self._add_textbox(slide, layout['date'], date)

    def add_section_break(self, title: str) -> None:
        """
        Add a section break/divider slide.

        Args:
            title: Section title text
        """
        self._slide_count += 1
        layout = LAYOUTS['section_break']

        slide = self.prs.slides.add_slide(self._get_blank_layout())
        self._set_background_color(slide, layout['background_color'])

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Logo mark (tree icon)
        if 'logo_mark' in layout:
            lm = layout['logo_mark']
            logo_mark_path = self.assets_dir / 'logo_mark.png'
            if logo_mark_path.exists():
                slide.shapes.add_picture(
                    str(logo_mark_path),
                    Inches(lm['x']), Inches(lm['y']),
                    Inches(lm['w']), Inches(lm['h'])
                )

        # Slide number
        if 'slide_number' in layout:
            sn = layout['slide_number']
            shape = slide.shapes.add_textbox(
                Inches(sn['x']), Inches(sn['y']),
                Inches(sn['w']), Inches(sn['h'])
            )
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = str(self._slide_count)
            p.font.name = sn.get('font', 'Avenir Medium')
            p.font.size = Pt(sn.get('size', 8))
            p.font.color.rgb = hex_to_rgb(sn.get('color', '#FFFFFF'))
            p.alignment = PP_ALIGN.RIGHT

    def add_content_slide(self, title: str, subtitle: str, bullets: list) -> None:
        """
        Add a text content slide with bullet points.

        Args:
            title: Slide title
            subtitle: Slide subtitle (not used in Stratfield template, kept for API compatibility)
            bullets: List of tuples (text, level) where level is 0-2
                    Example: [("Main point", 0), ("Sub point", 1)]
        """
        self._slide_count += 1
        layout = LAYOUTS['content_text_only']

        slide = self.prs.slides.add_slide(self._get_blank_layout())
        self._set_background_color(slide, layout['background_color'])

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Body with bullets
        body_spec = layout['body']
        body = slide.shapes.add_textbox(
            Inches(body_spec['x']),
            Inches(body_spec['y']),
            Inches(body_spec['w']),
            Inches(body_spec['h'])
        )

        tf = body.text_frame
        tf.word_wrap = True

        for i, (text, level) in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            self._apply_bullet_formatting(p, level, body_spec['color'])

        # Footer elements
        self._add_footer_elements(slide, layout)

    def add_image_slide(self, title: str, image_path: str, subtitle: str = "") -> None:
        """
        Add a slide with title and full-width image.

        Args:
            title: Slide title
            image_path: Path to image file
            subtitle: Optional subtitle (not used in Stratfield template, kept for API compatibility)
        """
        self._slide_count += 1
        layout = LAYOUTS['content_full_image']

        slide = self.prs.slides.add_slide(self._get_blank_layout())
        self._set_background_color(slide, layout['background_color'])

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Image
        self._add_image(slide, layout['image'], image_path)

        # Footer
        self._add_footer_elements(slide, layout)

    def add_text_and_image_slide(self, title: str, bullets: list, image_path: str) -> None:
        """
        Add a two-column slide with text on left and image on right.

        Args:
            title: Slide title (appears in dark left panel)
            bullets: List of tuples (text, level) for left panel
            image_path: Path to image for right side
        """
        self._slide_count += 1
        layout = LAYOUTS['content_text_and_image']

        slide = self.prs.slides.add_slide(self._get_blank_layout())
        self._set_background_color(slide, layout['background_color'])

        # Left panel (dark teal)
        self._add_rectangle(slide, layout['left_panel'])

        # Title (in left panel)
        self._add_textbox(slide, layout['title'], title)

        # Body with bullets (in left panel)
        body_spec = layout['body']
        body = slide.shapes.add_textbox(
            Inches(body_spec['x']),
            Inches(body_spec['y']),
            Inches(body_spec['w']),
            Inches(body_spec['h'])
        )

        tf = body.text_frame
        tf.word_wrap = True

        for i, (text, level) in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            # Use white bullets/text for dark panel background
            self._apply_bullet_formatting(p, level, body_spec['color'])
            # Override bullet color to white for visibility on dark background
            pPr = p._p.get_or_add_pPr()
            buClr = pPr.find(qn('a:buClr'))
            if buClr is not None:
                srgbClr = buClr.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    srgbClr.set('val', 'FFFFFF')

        # Image on right
        self._add_image(slide, layout['image'], image_path)

        # Footer
        self._add_footer_elements(slide, layout)

    def save(self, filepath: str) -> None:
        """
        Save the presentation to a file.

        Args:
            filepath: Output filename (should end in .pptx)
        """
        self.prs.save(filepath)
        print(f"Saved: {filepath}")

    def get_slide_count(self) -> int:
        """
        Get the current number of slides.

        Returns:
            Number of slides in the presentation
        """
        return self._slide_count
