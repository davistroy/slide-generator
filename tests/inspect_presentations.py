#!/usr/bin/env python3
"""
Detailed PowerPoint Inspection Script
Thoroughly inspects generated PPTX files and compares to source markdown.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

def inspect_slide(slide, slide_num):
    """Inspect a single slide and return detailed information."""
    info = {
        'number': slide_num,
        'shapes': [],
        'text_content': [],
        'has_title': False,
        'has_image': False,
        'has_table': False,
        'shape_count': len(slide.shapes)
    }

    for shape in slide.shapes:
        shape_info = {
            'type': shape.shape_type,
            'has_text': shape.has_text_frame if hasattr(shape, 'has_text_frame') else False,
        }

        # Check for title
        if hasattr(shape, 'name') and 'Title' in shape.name:
            info['has_title'] = True

        # Extract text
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text.strip()
            if text:
                info['text_content'].append(text)
                shape_info['text'] = text[:100] + '...' if len(text) > 100 else text

                # Count paragraphs/bullets
                if hasattr(shape.text_frame, 'paragraphs'):
                    shape_info['paragraph_count'] = len(shape.text_frame.paragraphs)

        # Check for tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            info['has_table'] = True
            shape_info['table_rows'] = len(shape.table.rows)
            shape_info['table_cols'] = len(shape.table.columns)

        # Check for pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info['has_image'] = True

        info['shapes'].append(shape_info)

    return info

def inspect_presentation(pptx_path, template_name):
    """Inspect entire presentation."""
    print(f"\n{'='*80}")
    print(f"INSPECTING: {pptx_path}")
    print(f"TEMPLATE: {template_name}")
    print(f"{'='*80}\n")

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"ERROR: Could not load presentation: {e}")
        return

    print(f"Total slides: {len(prs.slides)}")
    print(f"Slide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"Aspect ratio: {prs.slide_width.inches / prs.slide_height.inches:.2f}:1")
    print()

    # Inspect each slide
    issues = []

    for idx, slide in enumerate(prs.slides, 1):
        info = inspect_slide(slide, idx)

        print(f"--- SLIDE {idx} ---")
        print(f"  Shapes: {info['shape_count']}")
        print(f"  Has title: {info['has_title']}")
        print(f"  Has table: {info['has_table']}")
        print(f"  Has image: {info['has_image']}")

        # Show first few text items
        if info['text_content']:
            print(f"  Text preview:")
            for i, text in enumerate(info['text_content'][:3], 1):
                preview = text[:60] + '...' if len(text) > 60 else text
                print(f"    {i}. {preview}")
        else:
            print(f"  WARNING: No text content found!")
            issues.append(f"Slide {idx}: No text content")

        # Check for expected patterns
        if idx == 1:  # Title slide
            if not any('Block 1 Week 1' in text or 'Markdown' in text for text in info['text_content']):
                issues.append(f"Slide {idx}: Expected title slide content not found")

        if idx == 2:  # Should have table (Week Overview)
            if not info['has_table']:
                issues.append(f"Slide {idx}: Expected to have table but none found")

        print()

    # Summary
    print(f"\n{'='*80}")
    print(f"INSPECTION SUMMARY")
    print(f"{'='*80}\n")

    slides_with_text = sum(1 for slide in prs.slides if inspect_slide(slide, 0)['text_content'])
    slides_with_tables = sum(1 for slide in prs.slides if inspect_slide(slide, 0)['has_table'])
    slides_with_images = sum(1 for slide in prs.slides if inspect_slide(slide, 0)['has_image'])

    print(f"Slides with text: {slides_with_text}/{len(prs.slides)}")
    print(f"Slides with tables: {slides_with_tables}/{len(prs.slides)}")
    print(f"Slides with images: {slides_with_images}/{len(prs.slides)}")

    if issues:
        print(f"\nISSUES FOUND ({len(issues)}):")
        for issue in issues:
            print(f"  - {issue}")
    else:
        print(f"\n✓ No issues found!")

    return len(issues) == 0

if __name__ == "__main__":
    # Inspect both presentations
    cfa_ok = inspect_presentation("test-cfa.pptx", "CFA")
    stratfield_ok = inspect_presentation("test-stratfield.pptx", "Stratfield")

    print(f"\n{'='*80}")
    print(f"FINAL RESULTS")
    print(f"{'='*80}\n")
    print(f"CFA Presentation: {'PASS' if cfa_ok else 'FAIL'}")
    print(f"Stratfield Presentation: {'PASS' if stratfield_ok else 'FAIL'}")
    print()
