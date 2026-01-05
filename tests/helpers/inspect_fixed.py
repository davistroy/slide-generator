#!/usr/bin/env python3
"""
Inspect the fixed presentation and compare to expected content.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def inspect_slide_content(slide):
    """Extract all text content from a slide."""
    texts = []
    tables_count = 0

    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tables_count += 1

    return texts, tables_count


# Load the fixed presentation
prs = Presentation("test-cfa-fixed.pptx")

print("=" * 80)
print("FIXED PRESENTATION INSPECTION")
print("=" * 80)
print(f"\nTotal slides: {len(prs.slides)}\n")

# Check key slides
print("SLIDE 1 (TITLE SLIDE):")
texts, tables = inspect_slide_content(prs.slides[0])
print(f"  Texts found: {len(texts)}")
for i, text in enumerate(texts, 1):
    print(f"    {i}. {text}")
print(f"  Has ** markers: {'**' in ' '.join(texts)}")
print()

print("SLIDE 2 (WEEK OVERVIEW - should have table):")
texts, tables = inspect_slide_content(prs.slides[1])
print(f"  Texts found: {len(texts)}")
for i, text in enumerate(texts, 1):
    preview = text[:80] + "..." if len(text) > 80 else text
    print(f"    {i}. {preview}")
print(f"  Tables found: {tables}")
print(f"  Has ** markers: {'**' in ' '.join(texts)}")
print()

print("SLIDE 3 (MATURITY MODEL - should have table):")
texts, tables = inspect_slide_content(prs.slides[2])
print(f"  Texts found: {len(texts)}")
for i, text in enumerate(texts, 1):
    preview = text[:80] + "..." if len(text) > 80 else text
    print(f"    {i}. {preview}")
print(f"  Tables found: {tables}")
print(f"  Has ** markers: {'**' in ' '.join(texts)}")
print()

print("SLIDE 4 (LEARNING OBJECTIVES - numbered list):")
texts, tables = inspect_slide_content(prs.slides[3])
print(f"  Texts found: {len(texts)}")
for i, text in enumerate(texts[:5], 1):
    preview = text[:80] + "..." if len(text) > 80 else text
    print(f"    {i}. {preview}")
print(f"  Has ** markers: {'**' in ' '.join(texts)}")
print()

# Check all slides for ** markers
print("=" * 80)
print("CHECKING FOR MARKDOWN ARTIFACTS IN ALL SLIDES")
print("=" * 80)
slides_with_artifacts = []
for idx, slide in enumerate(prs.slides, 1):
    texts, _ = inspect_slide_content(slide)
    full_text = " ".join(texts)
    if "**" in full_text:
        slides_with_artifacts.append(idx)

if slides_with_artifacts:
    print(f"\n[FAIL] Found ** markers in slides: {slides_with_artifacts}")
else:
    print("\n[PASS] No ** markers found in any slide!")

# Count slides with content
slides_with_content = 0
empty_slides = []
for idx, slide in enumerate(prs.slides, 1):
    texts, _ = inspect_slide_content(slide)
    # Exclude slide numbers from content check
    content_texts = [t for t in texts if not t.isdigit()]
    if len(content_texts) >= 1:  # At least title
        slides_with_content += 1
    else:
        empty_slides.append(idx)

print(f"\n[INFO] Slides with content: {slides_with_content}/{len(prs.slides)}")
if empty_slides:
    print(f"[WARN] Empty slides: {empty_slides}")

print("\n" + "=" * 80)
print("SUMMARY")
print("=" * 80)
print(f"Total slides: {len(prs.slides)}")
print(f"Slides with content: {slides_with_content}")
print(f"Slides with ** artifacts: {len(slides_with_artifacts)}")
print(
    f"Tables found across all slides: {sum(inspect_slide_content(s)[1] for s in prs.slides)}"
)

if not slides_with_artifacts and slides_with_content == len(prs.slides):
    print("\n[SUCCESS] All critical issues fixed!")
else:
    print("\n[PARTIAL] Some issues remain")
