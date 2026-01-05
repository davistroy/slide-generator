from pptx import Presentation


def inspect(path, name):
    prs = Presentation(path)
    print(f"\n{'=' * 80}")
    print(f"Inspecting: {name}")
    print(f"{'=' * 80}")
    print(f"Total slides: {len(prs.slides)}")

    issues = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        images = 0
        tables = 0

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text
                text_parts.append(t)
                if "**" in t:
                    issues.append(f"Slide {slide_idx}: ** found")
                if "`" in t:
                    issues.append(f"Slide {slide_idx}: backtick found")

            if shape.shape_type == 13:
                images += 1
            if shape.shape_type == 19:
                tables += 1

        full = " ".join(text_parts)
        prev = full[:80] if len(full) > 80 else full
        print(
            f"  Slide {slide_idx}: {len(text_parts)} shapes, {images} images, {tables} tables | {prev}..."
        )

    print(f"\nIssues: {len(issues)}")
    for i in issues:
        print(f"  - {i}")
    return len(issues)


c = inspect("tests/artifacts/comprehensive-test-cfa.pptx", "CFA")
s = inspect("tests/artifacts/comprehensive-test-stratfield.pptx", "Stratfield")
print(f"\n{'=' * 80}")
print(f"TOTAL: CFA={c} issues, Stratfield={s} issues")
