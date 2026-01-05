#!/usr/bin/env python3
"""
Deep dive inspection of specific slides with issues.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def inspect_slide_deeply(prs, slide_num):
    """Deep inspection of a specific slide."""
    slide = prs.slides[slide_num - 1]

    print(f"\n{'=' * 80}")
    print(f"DEEP INSPECTION: SLIDE {slide_num}")
    print(f"{'=' * 80}\n")

    print(f"Total shapes: {len(slide.shapes)}\n")

    for i, shape in enumerate(slide.shapes, 1):
        print(f"--- Shape {i} ---")
        print(f"  Type: {shape.shape_type}")
        print(f"  Name: {shape.name if hasattr(shape, 'name') else 'N/A'}")

        if hasattr(shape, "text_frame"):
            print("  Has text frame: True")
            text = shape.text_frame.text
            print(f"  Text length: {len(text)}")
            print("  Text content:")
            print(f"    '''{text}'''")

            # Show paragraphs
            if hasattr(shape.text_frame, "paragraphs"):
                print(f"  Paragraphs: {len(shape.text_frame.paragraphs)}")
                for j, para in enumerate(shape.text_frame.paragraphs[:5], 1):
                    print(f"    Para {j}: level={para.level}, text='{para.text[:80]}'")

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print("  TABLE FOUND!")
            print(f"    Rows: {len(shape.table.rows)}")
            print(f"    Columns: {len(shape.table.columns)}")
            # Show table content
            for row_idx, row in enumerate(shape.table.rows[:3]):
                row_text = [cell.text_frame.text for cell in row.cells]
                print(f"    Row {row_idx}: {row_text}")

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print("  PICTURE FOUND")

        print()


# Inspect key slides
print("INSPECTING CFA PRESENTATION")
prs_cfa = Presentation("test-cfa.pptx")

# Slide 1 - Title slide
inspect_slide_deeply(prs_cfa, 1)

# Slide 2 - Week Overview (should have table)
inspect_slide_deeply(prs_cfa, 2)

# Slide 3 - Maturity Model (should have table)
inspect_slide_deeply(prs_cfa, 3)

# Slide 4 - Learning Objectives (bullets)
inspect_slide_deeply(prs_cfa, 4)

# Slide 5 - Problem Statement
inspect_slide_deeply(prs_cfa, 5)

print("\n" + "=" * 80)
print("CONTENT ISSUES SUMMARY")
print("=" * 80 + "\n")

# Check what the markdown actually contained for these slides
print("Expected content from markdown:")
print("\nSlide 2 (WEEK OVERVIEW) should have:")
print("  - Title: This Week's Journey")
print("  - Table with columns: Time | Topic | Focus")
print("  - 5 rows of session breakdown")

print("\nSlide 3 (MATURITY MODEL) should have:")
print("  - Title: Your Journey: From Ad-Hoc to Architect")
print("  - Table with columns: Level | Name | Description | Training")
print("  - 4 levels (0-3)")

print("\nSlide 4 (LEARNING OBJECTIVES) should have:")
print("  - Title: By the End of Today...")
print("  - Numbered list with 4 objectives")
print("  - Each with sub-bullets")
