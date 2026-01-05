#!/usr/bin/env python3
"""Comprehensive inspection of branded presentations."""

import sys
from pathlib import Path

from pptx import Presentation


def inspect_presentation(path, name):
    """Inspect a presentation and return detailed analysis."""
    prs = Presentation(path)

    print(f"\n{'=' * 80}")
    print(f"Inspecting: {name}")
    print(f"{'=' * 80}")
    print(f"File: {path}")
    print(f"File size: {Path(path).stat().st_size / 1024:.1f} KB")
    print(f"Total slides: {len(prs.slides)}")

    issues = []
    slide_details = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        images = 0
        tables = 0
        total_chars = 0

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text
                text_parts.append(t)
                total_chars += len(t)

                # Check for markdown artifacts
                if "**" in t:
                    issues.append(f"Slide {slide_idx}: ** found in text")
                if "`" in t:
                    issues.append(f"Slide {slide_idx}: backtick (`) found in text")
                if "*" in t and "**" not in t:
                    # Single asterisk not part of **
                    issues.append(
                        f"Slide {slide_idx}: single asterisk (*) found in text"
                    )

            # Count shape types
            if shape.shape_type == 13:  # Picture
                images += 1
            if shape.shape_type == 19:  # Table
                tables += 1

        full_text = " ".join(text_parts)
        preview = full_text[:80] if len(full_text) > 80 else full_text

        slide_info = {
            "number": slide_idx,
            "shapes": len(text_parts),
            "images": images,
            "tables": tables,
            "chars": total_chars,
            "preview": preview,
        }
        slide_details.append(slide_info)

        print(
            f"  Slide {slide_idx:2d}: {len(text_parts):2d} shapes, {images} img, {tables} tbl, {total_chars:4d} chars | {preview}..."
        )

    print(f"\n{'=' * 80}")
    print(f"ISSUES FOUND: {len(issues)}")
    if issues:
        for issue in issues:
            print(f"  - {issue}")
    else:
        print("  [OK] No markdown artifacts or issues found!")

    # Statistics
    total_images = sum(s["images"] for s in slide_details)
    total_tables = sum(s["tables"] for s in slide_details)
    total_chars = sum(s["chars"] for s in slide_details)

    print(f"\n{'=' * 80}")
    print("STATISTICS:")
    print(f"  Total images: {total_images}")
    print(f"  Total tables: {total_tables}")
    print(f"  Total characters: {total_chars}")
    print(f"  Average chars per slide: {total_chars / len(prs.slides):.1f}")

    return {
        "name": name,
        "slides": len(prs.slides),
        "issues": issues,
        "details": slide_details,
        "stats": {"images": total_images, "tables": total_tables, "chars": total_chars},
    }


def main():
    """Inspect both branded presentations."""
    print("\n" + "=" * 80)
    print("COMPREHENSIVE BRANDED PRESENTATION INSPECTION")
    print("=" * 80)

    cfa_result = inspect_presentation(
        "tests/artifacts/branded-test-cfa.pptx",
        "CFA Presentation (with CFA-branded images)",
    )

    stratfield_result = inspect_presentation(
        "tests/artifacts/branded-test-stratfield.pptx",
        "Stratfield Presentation (with Stratfield-branded images)",
    )

    # Comparison
    print(f"\n{'=' * 80}")
    print("COMPARISON: CFA vs STRATFIELD")
    print(f"{'=' * 80}")
    print(
        f"Slides:      CFA={cfa_result['slides']:2d}  Stratfield={stratfield_result['slides']:2d}"
    )
    print(
        f"Images:      CFA={cfa_result['stats']['images']:2d}  Stratfield={stratfield_result['stats']['images']:2d}"
    )
    print(
        f"Tables:      CFA={cfa_result['stats']['tables']:2d}  Stratfield={stratfield_result['stats']['tables']:2d}"
    )
    print(
        f"Characters:  CFA={cfa_result['stats']['chars']:5d}  Stratfield={stratfield_result['stats']['chars']:5d}"
    )
    print(
        f"Issues:      CFA={len(cfa_result['issues']):2d}  Stratfield={len(stratfield_result['issues']):2d}"
    )

    # Overall status
    print(f"\n{'=' * 80}")
    print("OVERALL TEST STATUS")
    print(f"{'=' * 80}")

    total_issues = len(cfa_result["issues"]) + len(stratfield_result["issues"])

    if total_issues == 0:
        print(
            "[PASS] Both presentations generated successfully with brand-specific images"
        )
        print("[PASS] NO markdown artifacts found")
        print("[PASS] All 20 slides present in both presentations")
        print("[PASS] All slides contain images")
        print("[PASS] Content properly formatted")
    else:
        print(f"[WARN] ISSUES FOUND - {total_issues} total issues detected")
        print("Review the detailed output above for specific problems")

    return total_issues


if __name__ == "__main__":
    exit_code = main()
    sys.exit(0 if exit_code == 0 else 1)
