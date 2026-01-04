"""
CFA PowerPoint Template v2.0

Chick-fil-A branded presentation template for the presentation-skill system.

Slide Types:
    - Title slide: Red background with white CFA compass logo
    - Section break: Dark blue (#004F71) background
    - Content slide: White background with bulleted content
    - Image slide: White background with full image
    - Text + Image slide: Gray left panel with bullets, image on right

Usage:
    from templates import get_template

    prs = get_template("cfa")
    prs.add_title_slide("Presentation Title", "Subtitle", "November 2025")
    prs.add_section_break("Section Name")
    prs.add_content_slide("Slide Title", "Subtitle", [
        ("Main bullet point", 0),
        ("Sub-bullet point", 1),
    ])
    prs.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
from typing import List, Tuple, Optional

from lib.template_base import PresentationTemplate
from templates import register_template


# ============================================================================
# BRAND COLORS
# ============================================================================
COLORS = {
    'cfa_red': '#DD0033',        # Primary CFA red
    'dark_blue': '#004F71',      # Section break background
    'light_gray': '#EEEDEB',     # Text+image panel background
    'text_gray': '#5B6770',      # Body text and bullets
    'white': '#FFFFFF',
    'teal': '#3EB1C8',           # Accent
    'green': '#249E6B',          # Accent
    'dark_teal': '#00635B',      # Accent
}

# ============================================================================
# FONTS
# ============================================================================
FONTS = {
    'heading': 'Apercu',
    'body': 'Apercu',
}

# ============================================================================
# SLIDE DIMENSIONS (Widescreen 16:9)
# ============================================================================
SLIDE_WIDTH = 13.33   # inches
SLIDE_HEIGHT = 7.50   # inches

# ============================================================================
# BULLET SPECIFICATIONS
# ============================================================================
BULLET_CHAR = '§'
BULLET_FONT = 'Wingdings'
BULLET_COLOR = '5B6770'  # Follows text color (gray)

BULLET_SPECS = {
    0: {
        'marL': 285750,      # 0.312"
        'indent': -285750,   # -0.312"
        'size': 15,
        'space_before': 6,
    },
    1: {
        'marL': 590550,      # 0.646"
        'indent': -285750,   # -0.312"
        'size': 13,
        'space_before': 4,
    },
    2: {
        'marL': 895350,      # 0.979"
        'indent': -285750,   # -0.312"
        'size': 11,
        'space_before': 4,
    },
}

# ============================================================================
# LAYOUT SPECIFICATIONS
# ============================================================================
LAYOUTS = {
    'title_slide': {
        'background_color': '#DD0033',
        'title': {
            'x': 0.31, 'y': 0.77, 'w': 10.64, 'h': 1.23,
            'font': 'Apercu', 'size': 44, 'color': '#FFFFFF',
            'align': 'left', 'bold': True
        },
        'subtitle': {
            'x': 0.38, 'y': 2.09, 'w': 6.0, 'h': 0.50,
            'font': 'Apercu', 'size': 18, 'color': '#FFFFFF',
            'align': 'left'
        },
        'date': {
            'x': 0.39, 'y': 6.40, 'w': 3.0, 'h': 0.30,
            'font': 'Apercu', 'size': 11, 'color': '#FFFFFF',
            'align': 'left'
        },
        'logo': {
            'x': 7.05, 'y': 2.10, 'w': 5.89, 'h': 5.00
        },
    },

    'section_break': {
        'background_color': '#004F71',
        'title': {
            'x': 1.12, 'y': 2.45, 'w': 10.98, 'h': 2.61,
            'font': 'Apercu', 'size': 44, 'color': '#FFFFFF',
            'align': 'left', 'bold': True
        },
        'chicken_icon': {
            'x': 0.11, 'y': 7.20, 'w': 0.17, 'h': 0.17
        },
        'slide_number': {
            'x': 12.84, 'y': 7.15, 'w': 0.40, 'h': 0.25,
            'font': 'Apercu', 'size': 9, 'color': '#FFFFFF',
            'align': 'right'
        },
    },

    'content_text_only': {
        'background_color': '#FFFFFF',
        'title': {
            'x': 0.06, 'y': 0.06, 'w': 12.54, 'h': 0.54,
            'font': 'Apercu', 'size': 24, 'color': '#5B6770',
            'align': 'left', 'bold': True
        },
        'subtitle': {
            'x': 0.06, 'y': 0.65, 'w': 12.54, 'h': 0.29,
            'font': 'Apercu', 'size': 12, 'color': '#5B6770',
            'align': 'left'
        },
        'body': {
            'x': 0.49, 'y': 1.20, 'w': 12.22, 'h': 5.69,
            'font': 'Apercu', 'color': '#5B6770',
        },
        'chicken_icon': {
            'x': 0.11, 'y': 7.20, 'w': 0.17, 'h': 0.17
        },
        'slide_number': {
            'x': 12.84, 'y': 7.15, 'w': 0.40, 'h': 0.25,
            'font': 'Apercu', 'size': 9, 'color': '#5B6770',
            'align': 'right'
        },
    },

    'content_full_image': {
        'background_color': '#FFFFFF',
        'title': {
            'x': 0.06, 'y': 0.06, 'w': 12.54, 'h': 0.54,
            'font': 'Apercu', 'size': 24, 'color': '#5B6770',
            'align': 'left', 'bold': True
        },
        'subtitle': {
            'x': 0.06, 'y': 0.65, 'w': 12.54, 'h': 0.29,
            'font': 'Apercu', 'size': 12, 'color': '#5B6770',
            'align': 'left'
        },
        'image': {
            'x': 0.58, 'y': 1.32, 'w': 12.00, 'h': 5.60,
        },
        'chicken_icon': {
            'x': 0.11, 'y': 7.20, 'w': 0.17, 'h': 0.17
        },
        'slide_number': {
            'x': 12.84, 'y': 7.15, 'w': 0.40, 'h': 0.25,
            'font': 'Apercu', 'size': 9, 'color': '#5B6770',
            'align': 'right'
        },
    },

    'content_text_and_image': {
        'background_color': '#FFFFFF',
        'left_panel': {
            'x': 0.0, 'y': 0.0, 'w': 5.36, 'h': 7.50,
            'fill': '#EEEDEB'
        },
        'title': {
            'x': 0.11, 'y': 0.16, 'w': 5.05, 'h': 0.78,
            'font': 'Apercu', 'size': 20, 'color': '#5B6770',
            'align': 'left', 'bold': True
        },
        'body': {
            'x': 0.11, 'y': 0.94, 'w': 5.05, 'h': 6.12,
            'font': 'Apercu', 'color': '#5B6770',
        },
        'image': {
            'x': 5.75, 'y': 0.37, 'w': 7.20, 'h': 6.76,
        },
        'chicken_icon': {
            'x': 0.11, 'y': 7.20, 'w': 0.17, 'h': 0.17
        },
        'slide_number': {
            'x': 12.84, 'y': 7.15, 'w': 0.40, 'h': 0.25,
            'font': 'Apercu', 'size': 9, 'color': '#5B6770',
            'align': 'right'
        },
    },
}


# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


# ============================================================================
# MAIN PRESENTATION CLASS
# ============================================================================

@register_template
class CFAPresentation(PresentationTemplate):
    """
    Chick-fil-A branded PowerPoint presentation template.

    Example:
        from templates import get_template

        prs = get_template("cfa")
        prs.add_title_slide("Q4 Business Review", "Operations Update", "November 2025")
        prs.add_section_break("Financial Overview")
        prs.add_content_slide("Revenue Performance", "Year-over-year comparison", [
            ("Total revenue increased 12% vs prior year", 0),
            ("Drive-thru sales up 15%", 1),
            ("Catering revenue up 8%", 1),
        ])
        prs.save("q4_review.pptx")
    """

    @property
    def name(self) -> str:
        """Template display name."""
        return "Chick-fil-A"

    @property
    def id(self) -> str:
        """Template identifier."""
        return "cfa"

    @property
    def description(self) -> str:
        """Template description."""
        return "Red and blue CFA branded template with Apercu font"

    def __init__(self, assets_dir: Optional[str] = None):
        """
        Initialize a new CFA presentation.

        Args:
            assets_dir: Path to assets folder. Defaults to 'assets' in same directory as this module.
        """
        self.prs = Presentation()

        # Set widescreen 16:9 dimensions
        self.prs.slide_width = Inches(SLIDE_WIDTH)
        self.prs.slide_height = Inches(SLIDE_HEIGHT)

        # Locate assets
        if assets_dir:
            self.assets_dir = Path(assets_dir)
        else:
            self.assets_dir = Path(__file__).parent / 'assets'

        self._slide_count = 0

    def _get_blank_layout(self):
        """Get or create a blank slide layout."""
        return self.prs.slide_layouts[6]  # Blank layout

    def _set_background_color(self, slide, color: str):
        """Set solid background color for slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color)

    def _add_textbox(self, slide, spec: dict, text: str, **overrides):
        """Add a formatted textbox to the slide."""
        merged = {**spec, **overrides}

        shape = slide.shapes.add_textbox(
            Inches(merged['x']),
            Inches(merged['y']),
            Inches(merged['w']),
            Inches(merged['h'])
        )

        tf = shape.text_frame
        tf.word_wrap = True

        # Set anchor if specified
        if 'anchor' in merged:
            anchor_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Note: vertical anchor set via shape properties

        p = tf.paragraphs[0]
        p.text = text

        # Alignment
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        p.alignment = align_map.get(merged.get('align', 'left'), PP_ALIGN.LEFT)

        # Font properties
        font = p.font
        font.name = merged.get('font', FONTS['body'])
        font.size = Pt(merged.get('size', 12))
        font.color.rgb = hex_to_rgb(merged.get('color', '#5B6770'))
        font.bold = merged.get('bold', False)

        return shape

    def _add_image(self, slide, spec: dict, image_path: str):
        """Add an image to the slide, fitting within specified bounds."""
        if not Path(image_path).exists():
            # Add placeholder rectangle if image doesn't exist
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(spec['x']), Inches(spec['y']),
                Inches(spec['w']), Inches(spec['h'])
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb('#CCCCCC')
            shape.line.fill.background()
            return shape

        # Calculate aspect-fit dimensions
        from PIL import Image
        with Image.open(image_path) as img:
            img_w, img_h = img.size

        img_aspect = img_w / img_h
        box_aspect = spec['w'] / spec['h']

        if img_aspect > box_aspect:
            # Image is wider - fit to width
            final_w = spec['w']
            final_h = spec['w'] / img_aspect
        else:
            # Image is taller - fit to height
            final_h = spec['h']
            final_w = spec['h'] * img_aspect

        # Center in bounding box
        x = spec['x'] + (spec['w'] - final_w) / 2
        y = spec['y'] + (spec['h'] - final_h) / 2

        return slide.shapes.add_picture(
            image_path,
            Inches(x), Inches(y),
            Inches(final_w), Inches(final_h)
        )

    def _apply_bullet_formatting(self, paragraph, level: int, bullet_color: str = BULLET_COLOR):
        """
        Apply PowerPoint native bullet formatting to a paragraph.

        Args:
            paragraph: The paragraph object to format
            level: Bullet level (0, 1, or 2)
            bullet_color: Hex color for bullet (without #)
        """
        spec = BULLET_SPECS.get(level, BULLET_SPECS[0])

        # Set text properties first (this creates defRPr element)
        paragraph.font.name = FONTS['body']
        paragraph.font.size = Pt(spec['size'])
        paragraph.font.color.rgb = hex_to_rgb('#' + BULLET_COLOR)
        paragraph.space_before = Pt(spec['space_before'])
        paragraph.level = level

        # Access the underlying XML element for bullet properties
        pPr = paragraph._p.get_or_add_pPr()

        # Set paragraph margins for proper bullet indentation
        pPr.set('marL', str(spec['marL']))
        pPr.set('indent', str(spec['indent']))

        # Remove any existing bullet settings
        for child in list(pPr):
            tag_name = etree.QName(child.tag).localname
            if tag_name in ('buNone', 'buChar', 'buFont', 'buClr', 'buAutoNum', 'buBlip', 'buClrTx'):
                pPr.remove(child)

        # Find defRPr to insert bullet elements BEFORE it
        # PowerPoint requires strict element ordering
        defRPr = pPr.find(qn('a:defRPr'))
        if defRPr is not None:
            insert_index = list(pPr).index(defRPr)
        else:
            insert_index = len(list(pPr))

        # Create bullet elements
        buClr = etree.Element(qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', bullet_color)

        buFont = etree.Element(qn('a:buFont'))
        buFont.set('typeface', BULLET_FONT)

        buChar = etree.Element(qn('a:buChar'))
        buChar.set('char', BULLET_CHAR)

        # Insert in correct order: buClr, buFont, buChar (before defRPr)
        pPr.insert(insert_index, buChar)
        pPr.insert(insert_index, buFont)
        pPr.insert(insert_index, buClr)

    def _add_footer_elements(self, slide, layout: dict, use_white_icon: bool = False):
        """Add chicken icon and slide number to slide."""
        # Chicken icon
        if 'chicken_icon' in layout:
            icon_spec = layout['chicken_icon']
            icon_name = 'chicken_white.png' if use_white_icon else 'chicken_red.png'
            icon_path = self.assets_dir / icon_name
            if icon_path.exists():
                slide.shapes.add_picture(
                    str(icon_path),
                    Inches(icon_spec['x']), Inches(icon_spec['y']),
                    Inches(icon_spec['w']), Inches(icon_spec['h'])
                )

        # Slide number
        if 'slide_number' in layout:
            sn = layout['slide_number']
            self._add_textbox(slide, sn, str(self._slide_count))

    # ========================================================================
    # PUBLIC METHODS - SLIDE TYPES (PresentationTemplate Interface)
    # ========================================================================

    def add_title_slide(self, title: str, subtitle: str = "", date: str = "") -> None:
        """
        Add a title/cover slide with red background and white CFA logo.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle
            date: Optional date or presenter info
        """
        self._slide_count += 1
        layout = LAYOUTS['title_slide']

        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Red background
        self._set_background_color(slide, layout['background_color'])

        # White CFA logo
        logo_path = self.assets_dir / 'logo_white.png'
        if logo_path.exists() and 'logo' in layout:
            logo_spec = layout['logo']
            slide.shapes.add_picture(
                str(logo_path),
                Inches(logo_spec['x']), Inches(logo_spec['y']),
                Inches(logo_spec['w']), Inches(logo_spec['h'])
            )

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Subtitle
        if subtitle:
            self._add_textbox(slide, layout['subtitle'], subtitle)

        # Date
        if date:
            self._add_textbox(slide, layout['date'], date)

    def add_section_break(self, title: str) -> None:
        """
        Add a section break slide with dark blue background.

        Args:
            title: Section title text
        """
        self._slide_count += 1
        layout = LAYOUTS['section_break']

        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Dark blue background
        self._set_background_color(slide, layout['background_color'])

        # Section title
        self._add_textbox(slide, layout['title'], title)

        # Footer with white icon (for dark background)
        self._add_footer_elements(slide, layout, use_white_icon=True)

    def add_content_slide(self, title: str, subtitle: str, bullets: List[Tuple[str, int]]) -> None:
        """
        Add a content slide with title, subtitle, and bulleted content.

        Args:
            title: Slide title
            subtitle: Slide subtitle (can be empty string)
            bullets: List of (text, level) tuples where level is 0, 1, or 2

        Example:
            prs.add_content_slide("Key Findings", "Q4 Analysis", [
                ("Revenue exceeded targets", 0),
                ("North region up 15%", 1),
                ("South region up 12%", 1),
                ("Cost management improved", 0),
            ])
        """
        self._slide_count += 1
        layout = LAYOUTS['content_text_only']

        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # White background (default)

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Subtitle
        if subtitle:
            self._add_textbox(slide, layout['subtitle'], subtitle)

        # Body with bullets
        body_spec = layout['body']
        body = slide.shapes.add_textbox(
            Inches(body_spec['x']),
            Inches(body_spec['y']),
            Inches(body_spec['w']),
            Inches(body_spec['h'])
        )

        tf = body.text_frame
        tf.word_wrap = True

        for i, (text, level) in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            self._apply_bullet_formatting(p, level)

        # Footer
        self._add_footer_elements(slide, layout)

    def add_image_slide(self, title: str, image_path: str, subtitle: str = "") -> None:
        """
        Add a slide with title and large image.

        Args:
            title: Slide title
            image_path: Path to image file
            subtitle: Optional subtitle
        """
        self._slide_count += 1
        layout = LAYOUTS['content_full_image']

        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Title
        self._add_textbox(slide, layout['title'], title)

        # Subtitle
        if subtitle:
            self._add_textbox(slide, layout['subtitle'], subtitle)

        # Image
        self._add_image(slide, layout['image'], image_path)

        # Footer
        self._add_footer_elements(slide, layout)

    def add_text_and_image_slide(self, title: str, bullets: List[Tuple[str, int]], image_path: str) -> None:
        """
        Add a slide with gray left panel (title + bullets) and image on right.

        Args:
            title: Title in left panel
            bullets: List of (text, level) tuples
            image_path: Path to image file
        """
        self._slide_count += 1
        layout = LAYOUTS['content_text_and_image']

        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Gray left panel
        panel_spec = layout['left_panel']
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(panel_spec['x']), Inches(panel_spec['y']),
            Inches(panel_spec['w']), Inches(panel_spec['h'])
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = hex_to_rgb(panel_spec['fill'])
        panel.line.fill.background()

        # Title in panel
        self._add_textbox(slide, layout['title'], title)

        # Body with bullets in panel
        body_spec = layout['body']
        body = slide.shapes.add_textbox(
            Inches(body_spec['x']),
            Inches(body_spec['y']),
            Inches(body_spec['w']),
            Inches(body_spec['h'])
        )

        tf = body.text_frame
        tf.word_wrap = True

        for i, (text, level) in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            self._apply_bullet_formatting(p, level)

        # Image on right
        self._add_image(slide, layout['image'], image_path)

        # Footer - use white icon on gray panel
        self._add_footer_elements(slide, layout, use_white_icon=False)

    def save(self, filepath: str) -> None:
        """
        Save the presentation to a file.

        Args:
            filepath: Output file path (should end in .pptx)
        """
        self.prs.save(filepath)
        print(f"Saved: {filepath}")

    def get_slide_count(self) -> int:
        """
        Get the current number of slides.

        Returns:
            Number of slides in the presentation
        """
        return self._slide_count

    def get_assets_dir(self) -> Path:
        """
        Get the path to this template's assets directory.

        Returns:
            Path to assets folder
        """
        return self.assets_dir
